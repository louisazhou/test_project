#!/usr/bin/env python3
"""
Clean slide system with decorator approach and figure generation functions.
"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from datetime import datetime
import pandas as pd
import numpy as np
from typing import Optional, Dict, Any, Callable, Tuple, List
from jinja2 import Template
import re
import matplotlib.pyplot as plt
import logging
from PIL import Image
import math

# Set up module logger
logger = logging.getLogger(__name__)

# Slide dimensions and margins (in inches)
SLIDE_WIDTH = 13.33  # 16:9 aspect ratio
SLIDE_HEIGHT = 7.5

# Margins and spacing
MARGIN_LEFT = 0.15   # Left margin 
MARGIN_RIGHT = 0.15  # Right margin 
MARGIN_TOP = 0.15    # Top margin 
MARGIN_BOTTOM = 0.15 # Bottom margin 

# Title area
TITLE_HEIGHT = 0.4   # Height for title 
TITLE_TOP = MARGIN_TOP
TITLE_LEFT = MARGIN_LEFT

# Content area calculations
CONTENT_LEFT = MARGIN_LEFT
CONTENT_TOP = TITLE_TOP + TITLE_HEIGHT + 0.1  # Small gap after title
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN_BOTTOM

# Spacing between elements
ELEMENT_GAP = 0.1    # Gap between elements 

# Text sizing constants
LINE_HEIGHT = 0.15   # Height per line of text (used for general layouts)
CHARS_PER_LINE = 120 # Characters per line at standard width

# Text estimation constants (for precise text height calculation)
TEXT_ESTIMATE_LINE_HEIGHT = 0.14      # Tighter line height for text estimation
TEXT_ESTIMATE_PARAGRAPH_SPACING = 0.04  # Minimal space between paragraphs  
TEXT_ESTIMATE_PADDING = 0.08          # Very conservative padding for text box

# Table sizing constants
TABLE_ROW_HEIGHT = 0.35  # Height per table row

# Default dimensions for figures and other elements
DEFAULT_FIGURE_HEIGHT = 4.0  # Standard figure height
MAX_CONTENT_HEIGHT = 6.0     # Maximum height for content before overflow
COLUMN_WIDTH_MULTIPLIER = 0.08  # Character length to inches conversion
MIN_COLUMN_WIDTH = 0.8       # Minimum column width in inches
MIN_INDEX_WIDTH = 1.3        # Minimum index column width in inches
MAX_TABLE_HEIGHT = 4.0       # Maximum table height in inches

# Layout validation constants
VALID_LAYOUT_TYPES = {
    'auto', 'text', 'text_figure', 'text_tables', 'text_tables_figure', 'text_tables_then_figure', 'text_figure_then_figure', 'figure'
}

class SlideContent:
    """Container for slide content with support for figure generation functions and multiple DataFrames."""
    
    def __init__(self, title: str, text_template: Optional[str] = None,
                 dfs: Optional[Dict[str, pd.DataFrame]] = None,
                 template_params: Optional[Dict] = None,
                 figure_generators: Optional[List[Dict]] = None,
                 show_figures: bool = True):
        self.title = title
        self.text_template = text_template
        self.dfs = dfs or {}  # Multiple DataFrames with marker names as keys
        self.template_params = template_params or {}
        self.figure_generators = figure_generators or []  # List of figure generator configs with functions
        self.show_figures = show_figures
        self._generated_figure_path = None
    
    def render_text(self) -> str:
        """Render the text using Jinja template, replacing table markers with placeholder text."""
        if not self.text_template:
            return ""
        
        # Create template params with table markers as placeholder text
        render_params = self.template_params.copy()
        
        # Add table markers as placeholder text for rendering
        for table_key in self.dfs.keys():
            render_params[table_key] = f"[TABLE: {table_key}]"
        
        template = Template(self.text_template)
        return template.render(**render_params)
    
    def get_text_and_table_positions(self) -> Tuple[List[str], List[Tuple[str, pd.DataFrame]]]:
        """
        Parse rendered text to extract text chunks and table positions.
        
        Returns:
            Tuple of (text_chunks, table_positions) where:
            - text_chunks: List of text segments
            - table_positions: List of (table_key, dataframe) tuples in order of appearance
        """
        rendered_text = self.render_text()
        
        if not self.dfs:
            return [rendered_text], []
        
        # Find all table markers in the rendered text
        text_chunks = []
        table_positions = []
        
        current_pos = 0
        
        # Find all table markers
        for match in re.finditer(r'\[TABLE: ([^\]]+)\]', rendered_text):
            # Add text before this table
            if match.start() > current_pos:
                text_before = rendered_text[current_pos:match.start()].strip()
                if text_before:
                    text_chunks.append(text_before)
            
            # Add table info
            table_key = match.group(1)
            if table_key in self.dfs:
                table_positions.append((table_key, self.dfs[table_key]))
            
            current_pos = match.end()
        
        # Add any remaining text after the last table
        if current_pos < len(rendered_text):
            remaining_text = rendered_text[current_pos:].strip()
            if remaining_text:
                text_chunks.append(remaining_text)
        
        # If no tables were found, return the full text
        if not table_positions:
            return [rendered_text], []
        
        return text_chunks, table_positions
    
    def get_all_dataframes(self) -> List[pd.DataFrame]:
        """Get all DataFrames in order of appearance in the template."""
        _, table_positions = self.get_text_and_table_positions()
        return [df for _, df in table_positions]
    
    def cleanup(self):
        """Clean up generated figure files."""
        if self._generated_figure_path and os.path.exists(self._generated_figure_path):
            os.unlink(self._generated_figure_path)
            self._generated_figure_path = None
    
    def render_console(self) -> str:
        """Render content for console display - shows text with embedded tables"""
        output = []
        
        # Get text chunks and table positions
        text_chunks, table_positions = self.get_text_and_table_positions()
        
        # Interleave text and tables
        chunk_idx = 0
        table_idx = 0
        
        if not table_positions:
            # No tables, just show text
            if text_chunks:
                output.append(text_chunks[0])
        else:
            # Interleave text and tables based on template structure
            # Simple approach: alternate between text chunks and tables
            while chunk_idx < len(text_chunks) or table_idx < len(table_positions):
                # Add text chunk if available
                if chunk_idx < len(text_chunks):
                    output.append(text_chunks[chunk_idx])
                    chunk_idx += 1
                
                # Add table if available
                if table_idx < len(table_positions):
                    table_key, df = table_positions[table_idx]
                    output.append(f"\n{table_key.upper()}:")
                    output.append(df.to_markdown())
                    table_idx += 1
        
        return '\n\n'.join(output)


class SlideLayouts:
    """Class for creating PowerPoint slides with automatic layout."""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)
    
    def _validate_layout_type(self, layout_type: str) -> str:
        """Validate and normalize layout type, raising error if invalid."""
        if layout_type not in VALID_LAYOUT_TYPES:
            raise ValueError(f"Invalid layout_type '{layout_type}'. Valid types are: {', '.join(sorted(VALID_LAYOUT_TYPES))}")
        return layout_type

    def _add_title(self, slide, title: str):
        """Add standardized title."""
        title_box = slide.shapes.add_textbox(
            Inches(TITLE_LEFT), 
            Inches(TITLE_TOP), 
            Inches(CONTENT_WIDTH),
            Inches(TITLE_HEIGHT)
        )
        title_box.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(18)
        title_box.text_frame.paragraphs[0].font.bold = True
        return title_box
    
    def _add_text(self, slide, text: str, left: float, top: float, 
                  width: float, height: float, font_size: int = 12):
        """Add text box with proper text wrapping and content-fitted sizing."""
        # Calculate the actual height needed for the text
        estimated_height = self._estimate_text_height(text)
        
        # Use the estimated height instead of the full available height
        actual_height = min(estimated_height, height)
        
        text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(actual_height))
        text_frame = text_box.text_frame
        
        # Enable word wrapping
        text_frame.word_wrap = True
        
        # Use text to fit shape for better content sizing
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Add text with proper paragraph formatting
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        
        # Set proper line spacing
        p.line_spacing = 1.0  # Single line spacing
        
        return text_box
    
    def _calculate_table_dimensions(self, df: pd.DataFrame, max_width: float = 7.0) -> tuple:
        """Calculate optimal table dimensions based on content length."""
        if df is None:
            return 0, 0
        
        include_index = not isinstance(df.index, pd.RangeIndex)
        
        # Calculate maximum content lengths
        if include_index:
            max_index_length = max(len(str(idx)) for idx in df.index)
            index_width = max(max_index_length * COLUMN_WIDTH_MULTIPLIER, MIN_INDEX_WIDTH)  # Use constant
        else:
            index_width = 0
            
        # Calculate data columns width
        max_col_lengths = []
        for col in df.columns:
            # Consider both header and content
            header_length = len(str(col))
            content_lengths = [len(str(val)) for val in df[col] if pd.notna(val)]
            max_length = max([header_length] + content_lengths) if content_lengths else header_length
            max_col_lengths.append(max_length)
        
        # Calculate total data width
        data_width = sum(max(length * COLUMN_WIDTH_MULTIPLIER, MIN_COLUMN_WIDTH) for length in max_col_lengths)
        
        # Add padding between columns (ELEMENT_GAP per column gap)
        padding = ELEMENT_GAP * (len(df.columns) - 1) if len(df.columns) > 1 else 0
        
        # Calculate total width with padding - use passed max_width parameter
        total_width = min(index_width + data_width + padding + 0.3, max_width)
        
        # Calculate height - simpler calculation since we're not wrapping content
        total_height = min((len(df) + 1) * TABLE_ROW_HEIGHT, MAX_TABLE_HEIGHT)  # +1 for header
        
        return total_width, total_height
    
    def _create_divider_slide(self, title: str) -> None:
        """Create a divider slide with just a title."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        title_box = self._add_title(slide, title)
        # Center the title vertically and make it larger
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.top = Inches(SLIDE_HEIGHT/2 - TITLE_HEIGHT/2)  # Center vertically
        return slide
    
    def _format_cell_value(self, value, is_percent=False):
        """Format cell value consistently."""
        if pd.isna(value):
            return "N/A"
        elif isinstance(value, int) or (isinstance(value, float) and value.is_integer()):
            # Handle integers (including floats that are actually integers)
            return str(int(value))
        elif isinstance(value, float):
            if is_percent:
                return f"{value*100:.0f}%"
            elif abs(value) < 0.01 and value != 0:
                return f"{value:.4f}"
            elif abs(value) < 1:
                return f"{value:.3f}"
            elif abs(value) < 100:
                return f"{value:.2f}"
            else:
                return f"{value:.0f}"
        return str(value)
    
    def _apply_cell_styling(self, cell, is_header=False, bg_color=None, text_color=None):
        """Apply consistent styling to a table cell."""
        # Set background color
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color or RGBColor(242, 244, 248)
        
        # Configure text frame
        text_frame = cell.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Style paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                if is_header:
                    run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = text_color or RGBColor(0, 0, 0)
    
    def _add_table(self, slide, df: pd.DataFrame, left: float, top: float, 
                   width: float, height: float, highlight_cells: Optional[Dict] = None):
        """Add DataFrame as table with intelligent column sizing."""
        highlight_cells = highlight_cells or {}
        
        include_index = not isinstance(df.index, pd.RangeIndex)
        rows, cols = df.shape[0] + 1, df.shape[1] + (1 if include_index else 0)
        
        # Use the passed width and height directly - they've already been calculated with constraints
        actual_width = width
        actual_height = height
        
        # Create table
        table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), 
                                           Inches(actual_width), Inches(actual_height))
        table = table_shape.table
        
        # Add headers
        for col_idx in range(cols):
            cell = table.cell(0, col_idx)
            if include_index and col_idx == 0:
                cell.text = ""
            else:
                actual_col_idx = col_idx - 1 if include_index else col_idx
                cell.text = str(df.columns[actual_col_idx])
            
            self._apply_cell_styling(
                cell, 
                is_header=True,
                bg_color=RGBColor(0, 0, 0),
                text_color=RGBColor(255, 255, 255)
            )
        
        # Add data
        for row_idx, (index_val, row) in enumerate(df.iterrows()):
            # Add index if needed
            if include_index:
                cell = table.cell(row_idx + 1, 0)
                cell.text = str(index_val)
                self._apply_cell_styling(
                    cell,
                    bg_color=RGBColor(242, 244, 248) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
                )
            
            # Add data cells
            for col_idx, col_name in enumerate(df.columns):
                cell = table.cell(row_idx + 1, col_idx + (1 if include_index else 0))
                
                # Get and format value
                value = df.iloc[row_idx, col_idx]
                # Enhanced percentage detection - check column name patterns and value ranges
                col_name_lower = str(col_name).lower()
                
                # Smart percentage detection: combine name-based and value-based detection
                # but exclude already-formatted data or percentage points
                name_indicates_percent = any(x in col_name_lower for x in ['pct', 'percentage', 'rate', 'ratio']) and not any(
                    x in col_name_lower for x in ['pp', 'delta', 'impact', 'δ']
                )
                value_indicates_percent = (
                    isinstance(value, (int, float)) and 
                    0 <= value <= 1 and 
                    value != 0 and
                    '(pp)' not in col_name_lower and  # Not percentage points
                    'δ' not in col_name_lower and     # Not deltas 
                    'impact' not in col_name_lower    # Not impact columns
                )
                
                is_percent = name_indicates_percent or value_indicates_percent
                cell.text = self._format_cell_value(value, is_percent)
                
                # Apply styling
                bg_color = RGBColor(242, 244, 248) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
                
                # Check for highlighting
                key = (index_val, col_name)
                if key in highlight_cells:
                    color = highlight_cells[key]
                    if color == "red":
                        text_color = RGBColor(255, 0, 0)
                    elif color == "green":
                        text_color = RGBColor(0, 153, 0)
                    elif isinstance(color, tuple) and len(color) == 3:
                        text_color = RGBColor(*color)
                    else:
                        text_color = RGBColor(0, 0, 0)
                else:
                    text_color = RGBColor(0, 0, 0)
                
                self._apply_cell_styling(cell, bg_color=bg_color, text_color=text_color)
        
        # Distribute column widths proportionally within the actual_width constraint
        if include_index:
            # For constrained tables, calculate proportional widths
            # Calculate what the natural index width would be
            max_index_length = max(len(str(idx)) for idx in df.index)
            natural_index_width = max(max_index_length * COLUMN_WIDTH_MULTIPLIER, MIN_INDEX_WIDTH)
            
            # Calculate natural data column widths
            natural_data_widths = []
            for col in df.columns:
                header_length = len(str(col))
                content_lengths = [len(str(val)) for val in df[col] if pd.notna(val)]
                max_length = max([header_length] + content_lengths) if content_lengths else header_length
                natural_width = max(max_length * COLUMN_WIDTH_MULTIPLIER, MIN_COLUMN_WIDTH)
                natural_data_widths.append(natural_width)
            
            total_natural_width = natural_index_width + sum(natural_data_widths)
            
            # If we need to constrain, scale all widths proportionally
            if total_natural_width > actual_width:
                scale_factor = actual_width / total_natural_width
                constrained_index_width = natural_index_width * scale_factor
                constrained_data_widths = [w * scale_factor for w in natural_data_widths]
            else:
                # Use natural widths if they fit
                constrained_index_width = natural_index_width
                constrained_data_widths = natural_data_widths
            
            # Apply the calculated widths
            for i in range(cols):
                if i == 0:
                    # Index column
                    table.columns[i].width = Inches(constrained_index_width)
                else:
                    # Data columns
                    table.columns[i].width = Inches(constrained_data_widths[i-1])
        else:
            # No index, distribute width equally among data columns
            col_width = actual_width / cols
            for i in range(cols):
                table.columns[i].width = Inches(col_width)
        
        return table
    
    def _add_caption(self, slide, caption: str, left: float, top: float, width: float, font_size: int = 10):
        """
        Add figure caption positioned below a figure.
        
        Args:
            slide: PowerPoint slide object
            caption: Caption text to display
            left, top: Position coordinates in inches
            width: Width of caption text box in inches
            font_size: Font size for caption text
        """
        if not caption.strip():
            return None
            
        # Estimate caption height (usually 1-2 lines)
        caption_height = max(0.3, len(caption) / 100 * 0.15)  # Rough estimate
        
        caption_box = slide.shapes.add_textbox(
            Inches(left), 
            Inches(top), 
            Inches(width),
            Inches(caption_height)
        )
        
        caption_frame = caption_box.text_frame
        caption_frame.word_wrap = True
        caption_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Add caption text with italics
        p = caption_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(font_size)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)  # Gray color
        p.alignment = PP_ALIGN.CENTER
        
        return caption_box
    
    def _add_figure(self, slide, figure_path: str, left: float, top: float, width: Optional[float] = None, height: Optional[float] = None, min_size: Optional[Dict[str, float]] = None, caption: Optional[str] = None):
        """
        Add figure to slide while preserving aspect ratio.
        Fits figure optimally within the specified width/height constraints.
        
        Args:
            slide: PowerPoint slide object
            figure_path: Path to the figure file
            left, top: Position coordinates in inches
            width, height: Optional size constraints in inches
            min_size: Optional dict with 'width' and 'height' minimum sizes in inches
            caption: Optional caption text to display below figure
        """
        if figure_path and os.path.exists(figure_path):
            # Get original image dimensions
            with Image.open(figure_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
            
            # Apply minimum size constraints if specified
            min_width = min_size.get('width', 0) if min_size else 0
            min_height = min_size.get('height', 0) if min_size else 0
            
            # Calculate dimensions that preserve aspect ratio
            if width is not None and height is not None:
                # Use the more constraining dimension
                width_from_height = height * aspect_ratio
                height_from_width = width / aspect_ratio
                
                if width_from_height <= width:
                    # Height is the constraining dimension
                    final_width = width_from_height
                    final_height = height
                else:
                    # Width is the constraining dimension
                    final_width = width
                    final_height = height_from_width
                    
                # Apply minimum size constraints
                if final_width < min_width:
                    final_width = min_width
                    final_height = min_width / aspect_ratio
                if final_height < min_height:
                    final_height = min_height
                    final_width = min_height * aspect_ratio
                    
            elif width is not None:
                # Only width specified
                final_width = max(width, min_width)
                final_height = max(final_width / aspect_ratio, min_height)
                # Recalculate width if height constraint was binding
                if final_height > final_width / aspect_ratio:
                    final_width = final_height * aspect_ratio
            elif height is not None:
                # Only height specified
                final_height = max(height, min_height)
                final_width = max(final_height * aspect_ratio, min_width)
                # Recalculate height if width constraint was binding
                if final_width > final_height * aspect_ratio:
                    final_height = final_width / aspect_ratio
            else:
                # No dimensions specified, use default size while preserving aspect ratio
                if aspect_ratio > 1:
                    # Landscape image
                    final_width = max(min(12.0, 8.0 * aspect_ratio), min_width)  # Max width 12 inches
                    final_height = max(final_width / aspect_ratio, min_height)
                else:
                    # Portrait image
                    final_height = max(min(6.0, 8.0 / aspect_ratio), min_height)  # Max height 6 inches
                    final_width = max(final_height * aspect_ratio, min_width)
            
            # Reserve space for caption if present
            caption_height = 0.0
            if caption and caption.strip():
                caption_height = max(0.3, len(caption) / 100 * 0.15) + 0.05  # Caption height + small gap
                
                # Adjust available height for figure if we have height constraint
                if height is not None:
                    available_height = height - caption_height
                    if available_height > 0:
                        # Recalculate figure dimensions with reduced height
                        height_from_width = final_width / aspect_ratio
                        if height_from_width > available_height:
                            final_height = available_height
                            final_width = available_height * aspect_ratio
                        else:
                            final_height = height_from_width
            
            # Add the figure with calculated dimensions
            slide.shapes.add_picture(
                figure_path, 
                Inches(left), 
                Inches(top),
                width=Inches(final_width),
                height=Inches(final_height)
            )
            
            # Add caption below the figure
            if caption and caption.strip():
                caption_top = top + final_height + 0.05  # Small gap between figure and caption
                self._add_caption(
                    slide=slide,
                    caption=caption,
                    left=left,
                    top=caption_top,
                    width=final_width
                )
        else:
            # Add placeholder text when no figure is available
            self._add_text(slide, f"Figure not available", left, top, 4.0, 1.0)
    
    def _estimate_text_height(self, text: str) -> float:
        """
        Estimate height needed for text with improved accuracy.
        Returns height in inches.
        """
        if not text:
            return 0.3
        
        # Clean the text first
        text = text.strip()
        if not text:
            return 0.3
        
        # Use tighter line height for text estimation
        
        # Split by paragraphs first (double newlines)
        paragraphs = text.split('\n\n')
        total_height = 0
        
        for para_idx, paragraph in enumerate(paragraphs):
            if not paragraph.strip():
                continue
                
            # Add paragraph spacing (except for first paragraph)
            if para_idx > 0:
                total_height += TEXT_ESTIMATE_PARAGRAPH_SPACING
            
            # Split paragraph into lines
            lines = paragraph.split('\n')
            paragraph_lines = 0
            
            for line in lines:
                if not line.strip():
                    paragraph_lines += 0.1  # Minimal space for empty lines
                    continue
                
                # Check for special formatting that needs extra space
                line_multiplier = 1.0
                if re.match(r'^\d+\.|\-|\*|\•', line.strip()):
                    line_multiplier = 1.02  # Tiny extra space for bullets/numbers
                elif line.strip().startswith('#') or line.strip().isupper():
                    line_multiplier = 1.05  # Small extra space for headers
                
                # Calculate wrapping more conservatively
                indent_level = len(re.match(r'^\s*', line).group())
                effective_width = CHARS_PER_LINE - (indent_level * 2)
                
                # Account for word boundaries - simplified approach
                line_length = len(line.strip())
                wrapped_lines = max(1, math.ceil(line_length / effective_width))
                
                paragraph_lines += wrapped_lines * line_multiplier
            
            total_height += paragraph_lines * TEXT_ESTIMATE_LINE_HEIGHT
        
        # Add padding
        total_height += TEXT_ESTIMATE_PADDING
        
        # Very minimal safety margin for PowerPoint quirks
        total_height *= 1.02
        
        return max(total_height, 0.3)  # Small minimum height
    
    def add_content_slide(self, content: SlideContent, layout_type: str = 'auto', 
                         highlight_cells: Optional[Dict] = None):
        """
        Add slide with automatic layout, content overflow detection, and multiple page support.
        
        Args:
            content: SlideContent object with all slide information
            layout_type: Layout preference ('auto', 'text_figure', 'text_tables', etc.)
            highlight_cells: Optional dictionary for table cell highlighting
        """
        # Validate layout type
        layout_type = self._validate_layout_type(layout_type)
        
        # Get text chunks and table positions
        text_chunks, table_positions = content.get_text_and_table_positions()
        
        # Handle figure generators - treat as overflow case if multiple figures
        figure_generators = content.figure_generators
        
        # Auto-detect layout if 'auto' is specified
        if layout_type == 'auto':
            has_text = len(text_chunks) > 0 and any(chunk.strip() for chunk in text_chunks)
            has_tables = len(table_positions) > 0
            has_figures = len(figure_generators) > 0
            
            if has_figures and has_text and has_tables:
                layout_type = 'text_tables_figure'
            elif has_figures and has_text:
                layout_type = 'text_figure'
            elif has_text and has_tables:
                layout_type = 'text_tables'
            elif has_figures:
                layout_type = 'figure'
            else:
                layout_type = 'text'  # Default fallback
        
        # Handle special two-slide layouts first
        if layout_type == 'text_tables_then_figure':
            self._create_text_tables_then_figure_slides(
                content=content,
                text_chunks=text_chunks,
                table_positions=table_positions,
                figure_generators=figure_generators,
                highlight_cells=highlight_cells
            )
            return
        elif layout_type == 'text_figure_then_figure':
            self._create_text_figure_then_figure_slides(
                content=content,
                text_chunks=text_chunks,
                table_positions=table_positions,
                figure_generators=figure_generators,
                highlight_cells=highlight_cells
            )
            return
        
        # Handle multiple figures from plotting functions that create figures in for-loops
        if len(figure_generators) > 1:
            # Multiple figures = overflow case, use enhanced overflow handling
            self._create_multiple_slides_for_overflow(
                content=content,
                text_chunks=text_chunks,
                table_positions=table_positions,
                figure_path=None,  # Figures handled within overflow method
                layout_type=layout_type,
                highlight_cells=highlight_cells,
                max_height=MAX_CONTENT_HEIGHT
            )
        elif len(figure_generators) == 1:
            # Single figure - check if it overflows with text/tables
            total_height = self._estimate_total_content_height(text_chunks, table_positions, "dummy_figure")
            max_height = MAX_CONTENT_HEIGHT
            
            if total_height > max_height:
                # Even single figure + content overflows
                self._create_multiple_slides_for_overflow(
                    content=content,
                    text_chunks=text_chunks,
                    table_positions=table_positions,
                    figure_path=None,  # Figure handled within overflow method
                    layout_type=layout_type,
                    highlight_cells=highlight_cells,
                    max_height=max_height
                )
            else:
                # Single figure + content fits on one slide
                self._create_single_figure_slide(content, text_chunks, table_positions, 
                                               figure_generators[0], layout_type, highlight_cells)
        else:
            # No figures - check for text/table overflow
            total_height = self._estimate_total_content_height(text_chunks, table_positions, None)
            max_height = MAX_CONTENT_HEIGHT
            
            # For text+tables content, check if two-column layout might fit better
            if (table_positions and self._should_use_two_column_layout(text_chunks, table_positions) 
                and total_height > max_height):
                
                # Estimate height for two-column layout (roughly half the height since content is split)
                estimated_two_column_height = total_height * 0.6  # Conservative estimate for two-column savings
                
                if estimated_two_column_height <= max_height:
                    # Two-column layout should fit - use it instead of overflow
                    self._create_slide_with_structured_content(
                        title=content.title,
                        text_chunks=text_chunks,
                        table_positions=table_positions,
                        figure_path=None,
                        layout_type=layout_type,
                        highlight_cells=highlight_cells
                    )
                    return
            
            # Handle different content types appropriately
            if total_height > max_height:
                if table_positions:
                    # Content with tables - use multi-slide approach for better table presentation
                    self._create_multiple_slides_for_overflow(
                        content=content,
                        text_chunks=text_chunks,
                        table_positions=table_positions,
                        figure_path=None,
                        layout_type=layout_type,
                        highlight_cells=highlight_cells,
                        max_height=max_height
                    )
                else:
                    # Text-only content - try truncation first, then multi-slide
                    self._create_slide_with_truncated_content(
                        content=content,
                        text_chunks=text_chunks,
                        table_positions=table_positions,
                        layout_type=layout_type,
                        highlight_cells=highlight_cells,
                        max_height=max_height
                    )
            else:
                # Content fits on single slide
                self._create_slide_with_structured_content(
                    title=content.title,
                    text_chunks=text_chunks,
                    table_positions=table_positions,
                    figure_path=None,
                    layout_type=layout_type,
                    highlight_cells=highlight_cells
                )
    
    def _create_single_figure_slide(self, content: SlideContent, text_chunks: List[str], 
                                   table_positions: List[Tuple[str, pd.DataFrame]], 
                                   fig_generator: Dict, layout_type: str, 
                                   highlight_cells: Optional[Dict]):
        """Create a single slide with one figure."""
        fig_params = fig_generator.get('params', {})
        fig_name = fig_generator.get('name', 'unknown')
        
        # Generate and add the figure
        generator_func = fig_generator.get('function')
        if generator_func:
            try:
                # Handle case where generator function creates multiple figures
                plt.close('all')  # Close any existing figures
                
                result = generator_func(**fig_params)
                
                # Handle different return types from generator functions
                figures_to_save = []
                if isinstance(result, plt.Figure):
                    figures_to_save.append(result)
                elif hasattr(result, '__iter__') and not isinstance(result, str):
                    # Handle case where function returns multiple figures
                    for item in result:
                        if isinstance(item, plt.Figure):
                            figures_to_save.append(item)
                
                # If no figures were returned directly, check all open figures
                if not figures_to_save:
                    figures_to_save = [plt.figure(num) for num in plt.get_fignums()]
                
                # Create slides for each figure
                for idx, fig in enumerate(figures_to_save):
                    # Create temp file
                    temp_dir = tempfile.gettempdir()
                    temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_{idx}.png"
                    figure_path = os.path.join(temp_dir, temp_filename)
                    
                    fig.savefig(figure_path, dpi=300, bbox_inches='tight')
                    
                    if content.show_figures:
                        plt.figure(fig.number)
                        plt.show()
                    
                    plt.close(fig)
                    
                    # Create slide title
                    slide_title = content.title
                    if len(figures_to_save) > 1:
                        if idx == 0:
                            slide_title = content.title  # First slide keeps original title
                        else:
                            slide_title = f"{content.title} (cont.)"  # Subsequent slides use (cont.)
                    elif fig_generator.get('title_suffix'):
                        slide_title = f"{content.title} - {fig_generator['title_suffix']}"
                    
                    # Only include text and tables on the first figure slide
                    slide_text_chunks = text_chunks if idx == 0 else []
                    slide_table_positions = table_positions if idx == 0 else []
                    
                    # Extract caption for proper handling (don't add to text chunks)
                    caption = fig_generator.get('caption', '')
                    
                    # Respect figure size - set minimum dimensions to prevent over-shrinking
                    min_figure_size = {'width': 4.0, 'height': 3.0}  # Minimum size in inches
                    
                    self._create_slide_with_structured_content(
                        title=slide_title,
                        text_chunks=slide_text_chunks,
                        table_positions=slide_table_positions,
                        figure_path=figure_path,
                        layout_type=layout_type,
                        highlight_cells=highlight_cells,
                        min_figure_size=min_figure_size,
                        figure_caption=caption if idx == 0 else None  # Only add caption to first slide
                    )
                    
                    # Clean up
                    if os.path.exists(figure_path):
                        os.unlink(figure_path)
                    
            except Exception as e:
                logger.warning(f"Failed to generate figure: {e}")
                self._create_slide_with_structured_content(
                    title=content.title,
                    text_chunks=text_chunks,
                    table_positions=table_positions,
                    figure_path=None,
                    layout_type=layout_type,
                    highlight_cells=highlight_cells
                )
        else:
            logger.warning(f"No generator function found in figure_generator config")
            self._create_slide_with_structured_content(
                title=content.title,
                text_chunks=text_chunks,
                table_positions=table_positions,
                figure_path=None,
                layout_type=layout_type,
                highlight_cells=highlight_cells
            )
    
    def _estimate_total_content_height(self, text_chunks: List[str], 
                                     table_positions: List[Tuple[str, pd.DataFrame]], 
                                     figure_path: Optional[str]) -> float:
        """Estimate total height needed for all content."""
        total_height = 0.0
        
        # Add text height with balanced spacing
        for i, text_chunk in enumerate(text_chunks):
            text_height = self._estimate_text_height(text_chunk)
            total_height += text_height
            
            # Add spacing between text chunks (except for last one)
            if i < len(text_chunks) - 1:
                total_height += 0.2  # Reasonable spacing between text sections
        
        # Add table heights with balanced spacing
        for i, (_, df) in enumerate(table_positions):
            _, table_height = self._calculate_table_dimensions(df)
            total_height += table_height
            
            # Add spacing between tables and after text
            if i == 0 and text_chunks:
                total_height += 0.25  # Reasonable space between text and first table
            elif i > 0:
                total_height += 0.2  # Space between tables
        
        # Add figure height
        if figure_path:
            total_height += DEFAULT_FIGURE_HEIGHT
            # Add spacing if there's other content
            if text_chunks or table_positions:
                total_height += 0.2
        
        return total_height
    
    def _create_multiple_slides_for_overflow(self, content: SlideContent, text_chunks: List[str],
                                           table_positions: List[Tuple[str, pd.DataFrame]],
                                           figure_path: Optional[str], layout_type: str,
                                           highlight_cells: Optional[Dict], max_height: float):
        """Create multiple slides when content overflows - handles text, tables, and figures."""
        slide_num = 1
        current_text = []
        current_height = 0.0
        
        # Get figure generators for handling
        figure_generators = content.figure_generators
        
        # Handle figures first - each figure gets its own slide
        for i, fig_generator in enumerate(figure_generators):
            # Generate and add the figure
            fig_params = fig_generator.get('params', {})
            generator_func = fig_generator.get('function')
            
            if generator_func:
                try:
                    plt.close('all')  # Close any existing figures
                    
                    result = generator_func(**fig_params)
                    
                    # Handle different return types from generator functions
                    figures_to_save = []
                    if isinstance(result, plt.Figure):
                        figures_to_save.append(result)
                    elif hasattr(result, '__iter__') and not isinstance(result, str):
                        # Handle case where function returns multiple figures
                        for item in result:
                            if isinstance(item, plt.Figure):
                                figures_to_save.append(item)
                    
                    # If no figures were returned directly, check all open figures
                    if not figures_to_save:
                        figures_to_save = [plt.figure(num) for num in plt.get_fignums()]
                    
                    # Create slides for each figure
                    for fig_idx, fig in enumerate(figures_to_save):
                        # Determine slide title
                        if fig_generator.get('title_suffix'):
                            if len(figures_to_save) > 1:
                                if fig_idx == 0:
                                    slide_title = f"{content.title} - {fig_generator['title_suffix']}"
                                else:
                                    slide_title = f"{content.title} - {fig_generator['title_suffix']} (cont.)"
                            else:
                                slide_title = f"{content.title} - {fig_generator['title_suffix']}"
                        else:
                            if slide_num == 1 and fig_idx == 0:
                                slide_title = content.title  # First slide keeps original title
                            else:
                                slide_title = f"{content.title} (cont.)"  # Use (cont.) for continuation
                        
                        # First figure slide gets text and tables, subsequent ones are figure-only
                        slide_text_chunks = text_chunks if i == 0 and fig_idx == 0 else []
                        slide_table_positions = table_positions if i == 0 and fig_idx == 0 else []
                        
                        # Extract caption for proper handling (don't add to text chunks)
                        caption = fig_generator.get('caption', '')
                        
                        # Create temp file
                        temp_dir = tempfile.gettempdir()
                        temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_{i}_{fig_idx}.png"
                        generated_figure_path = os.path.join(temp_dir, temp_filename)
                        
                        fig.savefig(generated_figure_path, dpi=300, bbox_inches='tight')
                        
                        if content.show_figures:
                            plt.figure(fig.number)
                            plt.show()
                        
                        plt.close(fig)
                        
                        # Respect figure size - set minimum dimensions
                        min_figure_size = {'width': 4.0, 'height': 3.0}
                        
                        # Create slide with figure
                        self._create_slide_with_structured_content(
                            title=slide_title,
                            text_chunks=slide_text_chunks,
                            table_positions=slide_table_positions,
                            figure_path=generated_figure_path,
                            layout_type=layout_type if slide_table_positions else ('text_figure' if slide_text_chunks else 'figure'),
                            highlight_cells=highlight_cells,
                            min_figure_size=min_figure_size,
                            figure_caption=caption if i == 0 and fig_idx == 0 else None  # Only add caption to first slide
                        )
                        
                        # Clean up
                        if os.path.exists(generated_figure_path):
                            os.unlink(generated_figure_path)
                        
                        slide_num += 1
                        
                except Exception as e:
                    logger.warning(f"Failed to generate figure: {e}")
                    slide_title = f"{content.title} (cont.)" if slide_num > 1 else content.title
                    slide_text_chunks = text_chunks if i == 0 else []
                    slide_table_positions = table_positions if i == 0 else []
                    
                    self._create_slide_with_structured_content(
                        title=slide_title,
                        text_chunks=slide_text_chunks,
                        table_positions=slide_table_positions,
                        figure_path=None,
                        layout_type='text_tables' if slide_text_chunks or slide_table_positions else 'text',
                        highlight_cells=highlight_cells
                    )
                    slide_num += 1
            else:
                logger.warning(f"No generator function found in figure_generator config")
                slide_title = f"{content.title} (cont.)" if slide_num > 1 else content.title
                slide_text_chunks = text_chunks if i == 0 else []
                slide_table_positions = table_positions if i == 0 else []
                
                self._create_slide_with_structured_content(
                    title=slide_title,
                    text_chunks=slide_text_chunks,
                    table_positions=slide_table_positions,
                    figure_path=None,
                    layout_type='text_tables' if slide_text_chunks or slide_table_positions else 'text',
                    highlight_cells=highlight_cells
                )
                slide_num += 1
        
        # If we had figures and text/tables were included in first slide, we're done
        if figure_generators and text_chunks:
            return
            
        # Process remaining text more intelligently - don't split aggressively
        full_text = '\n\n'.join(text_chunks)
        
        # First, check if the entire text actually fits (avoid unnecessary splitting)
        total_text_height = self._estimate_text_height(full_text)
        
        # Use 105% of max height instead of 95% to be more generous about fitting
        realistic_max_height = max_height * 1.05
        
        if total_text_height <= realistic_max_height:
            # All text fits on one slide - don't add (Part X) if it's the first slide!
            slide_title = content.title if slide_num == 1 else f"{content.title} (cont.)"
            self._create_slide_with_structured_content(
                title=slide_title,
                text_chunks=[full_text],
                table_positions=table_positions,
                figure_path=None,
                layout_type='text_tables' if table_positions else 'text',
                highlight_cells=highlight_cells
            )
            return
        
        # Only split when absolutely necessary - use larger meaningful sections
        # Split on double newlines (paragraphs) instead of numbered items
        paragraphs = full_text.split('\n\n')
        
        current_slide_text = []
        current_height = 0.0
        
        for paragraph in paragraphs:
            if not paragraph.strip():
                continue
                
            paragraph_height = self._estimate_text_height(paragraph)
            
            # Check if adding this paragraph would exceed reasonable slide capacity
            # Use 90% of max height for better readability and more conservative splitting
            if current_height + paragraph_height > max_height * 0.9 and current_slide_text:
                # Create slide with current content
                slide_title = content.title if slide_num == 1 else f"{content.title} (cont.)"
                self._create_slide_with_structured_content(
                    title=slide_title,
                    text_chunks=['\n\n'.join(current_slide_text)],
                    table_positions=[] if slide_num > 1 else table_positions,  # Tables only on first slide
                    figure_path=None,
                    layout_type='text_tables' if (slide_num == 1 and table_positions) else 'text',
                    highlight_cells=highlight_cells
                )
                slide_num += 1
                current_slide_text = []
                current_height = 0.0
            
            current_slide_text.append(paragraph)
            current_height += paragraph_height + 0.05  # Smaller gap between paragraphs
        
        # Create final slide with remaining content
        if current_slide_text:
            slide_title = content.title if slide_num == 1 else f"{content.title} (cont.)"
            self._create_slide_with_structured_content(
                title=slide_title,
                text_chunks=['\n\n'.join(current_slide_text)],
                table_positions=[] if slide_num > 1 else table_positions,  # Tables only if this is the first/only slide
                figure_path=None,
                layout_type='text_tables' if (slide_num == 1 and table_positions) else 'text',
                highlight_cells=highlight_cells
            )
        
        # Handle tables separately if they didn't fit on text slides
        if slide_num > 1 and table_positions:
            # Tables get their own slide(s)
            current_tables = []
            current_height = 0.0
            
            for table_key, df in table_positions:
                table_width, table_height = self._calculate_table_dimensions(df)
                
                if current_height + table_height > max_height * 0.8 and current_tables:
                    self._create_slide_with_structured_content(
                        title=f"{content.title} - Data Tables (cont.)",
                        text_chunks=[],
                        table_positions=current_tables,
                        figure_path=None,
                        layout_type='text_tables',
                        highlight_cells=highlight_cells
                    )
                    slide_num += 1
                    current_tables = []
                    current_height = 0.0
                
                current_tables.append((table_key, df))
                current_height += table_height + 0.2
            
            # Final table slide
            if current_tables:
                self._create_slide_with_structured_content(
                    title=f"{content.title} - Data Tables (cont.)",
                    text_chunks=[],
                    table_positions=current_tables,
                    figure_path=None,
                    layout_type='text_tables',
                    highlight_cells=highlight_cells
                )
    
    def _create_slide_with_structured_content(self, title: str, text_chunks: List[str], 
                                            table_positions: List[Tuple[str, pd.DataFrame]], 
                                            figure_path: Optional[str], layout_type: str, 
                                            highlight_cells: Optional[Dict],
                                            min_figure_size: Optional[Dict[str, float]] = None,
                                            figure_caption: Optional[str] = None):
        """Create slide with structured content positioning using helper functions."""
        # Validate layout type
        layout_type = self._validate_layout_type(layout_type)
        
        # Create basic slide structure
        slide = self._create_basic_slide_structure(title)
        
        # Check if we should use intelligent two-column layout for text-tables
        if layout_type == 'text_tables' and self._should_use_two_column_layout(text_chunks, table_positions):
            self._create_two_column_text_tables_layout(slide, text_chunks, table_positions, highlight_cells)
            return slide
        
        # Standard single-column layouts
        # Combine text chunks
        full_text = self._combine_text_chunks(text_chunks)
        text_height = self._estimate_text_height(full_text) if full_text else 0
        
        # Calculate table dimensions if tables exist
        table_width = table_height = 0
        if table_positions:
            table_width, table_height = self._calculate_table_dimensions(table_positions[0][1])
        
        # Calculate layout positions
        layouts = self._calculate_content_layout(
            content_type=layout_type,
            text_height=text_height,
            table_width=table_width,
            table_height=table_height
        )
        
        # Add content blocks based on layout type
        if layout_type == 'text' and full_text:
            self._add_content_block(
                slide=slide,
                content_type='text',
                content_data={'text': full_text},
                layout_info=layouts['text']
            )
            
        elif layout_type == 'text_figure':
            if full_text:
                self._add_content_block(
                    slide=slide,
                    content_type='text',
                    content_data={'text': full_text},
                    layout_info=layouts['text']
                )
            if figure_path:
                self._add_content_block(
                    slide=slide,
                    content_type='figure',
                    content_data={'figure_path': figure_path, 'min_size': min_figure_size, 'caption': figure_caption},
                    layout_info=layouts['figure']
                )
                
        elif layout_type == 'figure' and figure_path:
            figure_layout = {
                'left': CONTENT_LEFT,
                'top': CONTENT_TOP,
                'width': CONTENT_WIDTH,
                'height': CONTENT_HEIGHT
            }
            self._add_content_block(
                slide=slide,
                content_type='figure',
                content_data={'figure_path': figure_path, 'min_size': min_figure_size, 'caption': figure_caption},
                layout_info=figure_layout
            )
            
        elif layout_type == 'text_tables':
            if full_text:
                self._add_content_block(
                    slide=slide,
                    content_type='text',
                    content_data={'text': full_text},
                    layout_info=layouts['text']
                )
            
            # Add tables sequentially
            current_top = layouts['tables']['top']
            for table_key, df in table_positions:
                table_width, table_height = self._calculate_table_dimensions(df)
                table_layout = {
                    'left': layouts['tables']['left'],
                    'top': current_top,
                    'width': table_width,
                    'height': table_height
                }
                self._add_content_block(
                    slide=slide,
                    content_type='table',
                    content_data={'df': df, 'highlight_cells': highlight_cells},
                    layout_info=table_layout
                )
                current_top += table_height + ELEMENT_GAP
                
        elif layout_type == 'text_tables_figure':
            if full_text:
                self._add_content_block(
                    slide=slide,
                    content_type='text',
                    content_data={'text': full_text},
                    layout_info=layouts['text']
                )
            
            if table_positions:
                table_key, df = table_positions[0]
                self._add_content_block(
                    slide=slide,
                    content_type='table',
                    content_data={'df': df, 'highlight_cells': highlight_cells},
                    layout_info=layouts['tables']
                )
            
            if figure_path:
                self._add_content_block(
                    slide=slide,
                    content_type='figure',
                    content_data={'figure_path': figure_path, 'min_size': min_figure_size, 'caption': figure_caption},
                    layout_info=layouts['figure']
                )
        
        return slide
    
    def _create_text_tables_then_figure_slides(self, content: SlideContent, text_chunks: List[str],
                                              table_positions: List[Tuple[str, pd.DataFrame]],
                                              figure_generators: List[Dict],
                                              highlight_cells: Optional[Dict]):
        """Create two slides: first with text+tables, second with figure+caption."""
        
        # Slide 1: Text + Tables (full width)
        self._create_slide_with_structured_content(
            title=content.title,
            text_chunks=text_chunks,
            table_positions=table_positions,
            figure_path=None,
            layout_type='text_tables',
            highlight_cells=highlight_cells
        )
        
        # Slide 2: Figure + Caption (if we have figures)
        if figure_generators:
            for i, fig_generator in enumerate(figure_generators):
                fig_params = fig_generator.get('params', {})
                generator_func = fig_generator.get('function')
                
                if generator_func:
                    try:
                        plt.close('all')  # Close any existing figures
                        
                        result = generator_func(**fig_params)
                        
                        # Handle different return types from generator functions
                        figures_to_save = []
                        if isinstance(result, plt.Figure):
                            figures_to_save.append(result)
                        elif hasattr(result, '__iter__') and not isinstance(result, str):
                            # Handle case where function returns multiple figures
                            for item in result:
                                if isinstance(item, plt.Figure):
                                    figures_to_save.append(item)
                        
                        # If no figures were returned directly, check all open figures
                        if not figures_to_save:
                            figures_to_save = [plt.figure(num) for num in plt.get_fignums()]
                        
                        # Create slides for each figure
                        for fig_idx, fig in enumerate(figures_to_save):
                            # Create temp file
                            temp_dir = tempfile.gettempdir()
                            temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_{i}_{fig_idx}.png"
                            figure_path = os.path.join(temp_dir, temp_filename)
                            
                            fig.savefig(figure_path, dpi=300, bbox_inches='tight')
                            
                            if content.show_figures:
                                plt.figure(fig.number)
                                plt.show()
                            
                            plt.close(fig)
                            
                            # Create slide title for figure slide
                            if fig_generator.get('title_suffix'):
                                slide_title = f"{content.title} - {fig_generator['title_suffix']}"
                            else:
                                slide_title = content.title
                            
                            # Extract caption
                            caption = fig_generator.get('caption', '')
                            
                            # Create figure slide (no text, no tables)
                            self._create_slide_with_structured_content(
                                title=slide_title,
                                text_chunks=[],
                                table_positions=[],
                                figure_path=figure_path,
                                layout_type='figure',
                                highlight_cells=highlight_cells,
                                figure_caption=caption
                            )
                            
                            # Clean up
                            if os.path.exists(figure_path):
                                os.unlink(figure_path)
                            
                    except Exception as e:
                        logger.warning(f"Failed to generate figure: {e}")
                        # Create a placeholder slide if figure generation fails
                        self._create_slide_with_structured_content(
                            title=content.title,
                            text_chunks=["Figure generation failed."],
                            table_positions=[],
                            figure_path=None,
                            layout_type='text',
                            highlight_cells=highlight_cells
                        )
    
    def _create_text_figure_then_figure_slides(self, content: SlideContent, text_chunks: List[str],
                                              table_positions: List[Tuple[str, pd.DataFrame]],
                                              figure_generators: List[Dict],
                                              highlight_cells: Optional[Dict]):
        """Create two slides: first with text+first figure, second with subsequent figures."""
        
        if not figure_generators:
            # No figures, just create text slide
            self._create_slide_with_structured_content(
                title=content.title,
                text_chunks=text_chunks,
                table_positions=table_positions,
                figure_path=None,
                layout_type='text_tables' if table_positions else 'text',
                highlight_cells=highlight_cells
            )
            return
        
        # Slide 1: Text + First Figure
        first_fig_generator = figure_generators[0]
        fig_params = first_fig_generator.get('params', {})
        generator_func = first_fig_generator.get('function')
        
        if generator_func:
            try:
                plt.close('all')  # Close any existing figures
                
                result = generator_func(**fig_params)
                
                # Handle different return types from generator functions
                figures_to_save = []
                if isinstance(result, plt.Figure):
                    figures_to_save.append(result)
                elif hasattr(result, '__iter__') and not isinstance(result, str):
                    # Handle case where function returns multiple figures
                    for item in result:
                        if isinstance(item, plt.Figure):
                            figures_to_save.append(item)
                
                # If no figures were returned directly, check all open figures
                if not figures_to_save:
                    figures_to_save = [plt.figure(num) for num in plt.get_fignums()]
                
                # Use first figure for the text+figure slide
                if figures_to_save:
                    first_fig = figures_to_save[0]
                    
                    # Create temp file for first figure
                    temp_dir = tempfile.gettempdir()
                    temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_first.png"
                    first_figure_path = os.path.join(temp_dir, temp_filename)
                    
                    first_fig.savefig(first_figure_path, dpi=300, bbox_inches='tight')
                    
                    if content.show_figures:
                        plt.figure(first_fig.number)
                        plt.show()
                    
                    plt.close(first_fig)
                    
                    # Extract caption
                    caption = first_fig_generator.get('caption', '')
                    
                    # Create first slide with text + first figure
                    self._create_slide_with_structured_content(
                        title=content.title,
                        text_chunks=text_chunks,
                        table_positions=table_positions,
                        figure_path=first_figure_path,
                        layout_type='text_tables_figure' if table_positions else 'text_figure',
                        highlight_cells=highlight_cells,
                        figure_caption=caption
                    )
                    
                    # Clean up first figure
                    if os.path.exists(first_figure_path):
                        os.unlink(first_figure_path)
                    
                    # Handle remaining figures from the first generator
                    for fig_idx, fig in enumerate(figures_to_save[1:], 1):
                        temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_0_{fig_idx}.png"
                        figure_path = os.path.join(temp_dir, temp_filename)
                        
                        fig.savefig(figure_path, dpi=300, bbox_inches='tight')
                        
                        if content.show_figures:
                            plt.figure(fig.number)
                            plt.show()
                        
                        plt.close(fig)
                        
                        # Create figure slide (no text, no tables)
                        self._create_slide_with_structured_content(
                            title=content.title,
                            text_chunks=[],
                            table_positions=[],
                            figure_path=figure_path,
                            layout_type='figure',
                            highlight_cells=highlight_cells,
                            figure_caption=caption
                        )
                        
                        # Clean up
                        if os.path.exists(figure_path):
                            os.unlink(figure_path)
                
            except Exception as e:
                logger.warning(f"Failed to generate first figure: {e}")
                # Create text-only slide if figure generation fails
                self._create_slide_with_structured_content(
                    title=content.title,
                    text_chunks=text_chunks,
                    table_positions=table_positions,
                    figure_path=None,
                    layout_type='text_tables' if table_positions else 'text',
                    highlight_cells=highlight_cells
                )
        
        # Slide 2+: Subsequent Figures (if any)
        for i, fig_generator in enumerate(figure_generators[1:], 1):
            fig_params = fig_generator.get('params', {})
            generator_func = fig_generator.get('function')
            
            if generator_func:
                try:
                    plt.close('all')  # Close any existing figures
                    
                    result = generator_func(**fig_params)
                    
                    # Handle different return types from generator functions
                    figures_to_save = []
                    if isinstance(result, plt.Figure):
                        figures_to_save.append(result)
                    elif hasattr(result, '__iter__') and not isinstance(result, str):
                        # Handle case where function returns multiple figures
                        for item in result:
                            if isinstance(item, plt.Figure):
                                figures_to_save.append(item)
                    
                    # If no figures were returned directly, check all open figures
                    if not figures_to_save:
                        figures_to_save = [plt.figure(num) for num in plt.get_fignums()]
                    
                    # Create slides for each figure
                    for fig_idx, fig in enumerate(figures_to_save):
                        # Create temp file
                        temp_dir = tempfile.gettempdir()
                        temp_filename = f"temp_figure_{int(datetime.now().timestamp() * 1000)}_{i}_{fig_idx}.png"
                        figure_path = os.path.join(temp_dir, temp_filename)
                        
                        fig.savefig(figure_path, dpi=300, bbox_inches='tight')
                        
                        if content.show_figures:
                            plt.figure(fig.number)
                            plt.show()
                        
                        plt.close(fig)
                        
                        # Extract caption
                        caption = fig_generator.get('caption', '')
                        
                        # Create figure slide (no text, no tables)
                        self._create_slide_with_structured_content(
                            title=content.title,
                            text_chunks=[],
                            table_positions=[],
                            figure_path=figure_path,
                            layout_type='figure',
                            highlight_cells=highlight_cells,
                            figure_caption=caption
                        )
                        
                        # Clean up
                        if os.path.exists(figure_path):
                            os.unlink(figure_path)
                        
                except Exception as e:
                    logger.warning(f"Failed to generate figure: {e}")
                    # Create a placeholder slide if figure generation fails
                    self._create_slide_with_structured_content(
                        title=content.title,
                        text_chunks=["Figure generation failed."],
                        table_positions=[],
                        figure_path=None,
                        layout_type='text',
                        highlight_cells=highlight_cells
                    )
    
    def _create_slide_with_truncated_content(self, content: SlideContent, text_chunks: List[str],
                                           table_positions: List[Tuple[str, pd.DataFrame]],
                                           layout_type: str, highlight_cells: Optional[Dict],
                                           max_height: float):
        """Create a single slide with intelligently truncated text content."""
        # Combine all text chunks
        full_text = '\n\n'.join(text_chunks)
        
        # Calculate available height for text (accounting for title and margins)
        available_text_height = max_height - 0.5  # Reserve space for title and margins
        
        # Truncate text to fit
        truncated_text = self._truncate_text_to_fit(full_text, available_text_height)
        
        # Create slide with truncated content
        self._create_slide_with_structured_content(
            title=content.title,
            text_chunks=[truncated_text],
            table_positions=table_positions,
            figure_path=None,
            layout_type=layout_type,
            highlight_cells=highlight_cells
        )
    
    def _truncate_text_to_fit(self, text: str, max_height: float) -> str:
        """
        Intelligently truncate text to fit within the specified height.
        Truncates at paragraph boundaries and adds ellipsis if truncated.
        """
        if not text:
            return text
        
        # Split text into paragraphs
        paragraphs = text.split('\n\n')
        result_paragraphs = []
        current_height = 0.0
        
        for paragraph in paragraphs:
            # Estimate height of this paragraph
            paragraph_height = self._estimate_text_height(paragraph.strip())
            
            # Check if adding this paragraph would exceed max height
            if current_height + paragraph_height > max_height:
                # If we haven't added any paragraphs yet, try to fit part of this one
                if not result_paragraphs:
                    # Split by sentences and try to fit what we can
                    sentences = paragraph.split('. ')
                    for i, sentence in enumerate(sentences):
                        partial_text = '. '.join(sentences[:i+1])
                        if not partial_text.endswith('.'):
                            partial_text += '.'
                        
                        partial_height = self._estimate_text_height(partial_text)
                        if partial_height > max_height:
                            # Even this sentence doesn't fit, take what we can
                            if i > 0:
                                result_paragraphs.append('. '.join(sentences[:i]) + '.')
                            break
                        elif i == len(sentences) - 1:
                            # All sentences fit
                            result_paragraphs.append(partial_text)
                        # Continue to next sentence
                
                # Add truncation indicator and break
                if result_paragraphs:
                    result_text = '\n\n'.join(result_paragraphs)
                    # Add ellipsis if we truncated
                    if len(result_text.strip()) < len(text.strip()):
                        result_text += "\n\n[Content truncated - see full analysis for complete details]"
                    return result_text
                else:
                    # Couldn't fit anything, return a summary message
                    return "[Content too large for slide - please see full document]"
            
            # This paragraph fits, add it
            result_paragraphs.append(paragraph)
            current_height += paragraph_height + 0.1  # Add small gap between paragraphs
        
        # All paragraphs fit
        return '\n\n'.join(result_paragraphs)
    
    def save(self, filename: Optional[str] = None, output_dir: str = '.'):
        """Save presentation."""
        os.makedirs(output_dir, exist_ok=True)
        
        if not filename:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"presentation_{timestamp}.pptx"
        elif not filename.endswith('.pptx'):
            filename = f"{filename}.pptx"
        
        filepath = os.path.join(output_dir, filename)
        self.prs.save(filepath)
        return filepath
    
    def _create_highlight_cells_map(self, df: pd.DataFrame, metric_anomaly_map: Dict[str, Dict[str, Any]]) -> Dict[Tuple[str, str], str]:
        """Create a highlight_cells dictionary for the table based on anomaly information."""
        highlight_cells = {}
        
        for metric_name, metric_info in metric_anomaly_map.items():
            if metric_name in df.index:
                anomalous_region = metric_info.get('anomalous_region')
                direction = metric_info.get('direction')
                higher_is_better = metric_info.get('higher_is_better', True)
                
                if anomalous_region and direction and anomalous_region in df.columns:
                    is_good = (direction == 'higher' and higher_is_better) or \
                             (direction == 'lower' and not higher_is_better)
                    
                    color = 'green' if is_good else 'red'
                    highlight_cells[(metric_name, anomalous_region)] = color
        
        return highlight_cells
    
    def create_metrics_summary_slide(self, df: pd.DataFrame, metrics_text: Dict[str, str], 
                                   metric_anomaly_map: Dict[str, Dict[str, Any]], 
                                   title: str = "Metrics Summary", 
                                   max_table_width: Optional[float] = None) -> None:
        """
        Create a summary slide with a styled table of metrics data and pre-formatted text explanations.
        
        Args:
            df: DataFrame with metrics data
            metrics_text: Dictionary mapping metric names to text explanations
            metric_anomaly_map: Dictionary with anomaly information for highlighting
            title: Slide title
            max_table_width: Maximum width for the table (useful for narrow tables to prevent text overlap)
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add standardized title using helper
        self._add_title(slide, title)
        
        # Filter DataFrame to only include metrics we have text for
        # metrics = list(metrics_text.keys())
        # metrics_df = df[metrics]
        metrics_df = df
        
        # Transpose the DataFrame so regions are columns and metrics are rows
        metrics_df_transposed = metrics_df.T
        
        # Create highlight cells map for transposed data
        highlight_cells = self._create_highlight_cells_map(metrics_df_transposed, metric_anomaly_map)
        
        # Calculate optimal table dimensions with optional width constraint
        # Use max_table_width directly if specified, otherwise use default
        max_width_for_calculation = max_table_width if max_table_width is not None else 6.0
        calculated_table_width, table_height = self._calculate_table_dimensions(metrics_df_transposed, max_width=max_width_for_calculation)
        
        # The calculated width already respects max_table_width, so use it directly
        table_width = calculated_table_width + 0.2  # Safety margin for PowerPoint spacing
        
        # Position table on the left side using standard constants
        table_left = CONTENT_LEFT
        table_top = CONTENT_TOP
        
        # Add styled table using helper method - table_width here controls actual table rendering
        self._add_table(
            slide=slide,
            df=metrics_df_transposed,
            left=table_left,
            top=table_top,
            width=table_width,
            height=table_height,
            highlight_cells=highlight_cells
        )
        
        # Calculate text area on the right side with minimal spacing
        text_left = table_left + table_width + ELEMENT_GAP + 0.03  # Reduced spacing
        # Calculate available width from text start to slide end
        text_width = (CONTENT_LEFT + CONTENT_WIDTH) - text_left
        
        # Ensure text doesn't go beyond slide boundaries (already handled above)
        if text_width < 0:
            text_width = 0
        
        # Ensure minimum text width
        if text_width < 2.0:
            text_width = 2.0
        
        text_height = CONTENT_HEIGHT
        
        # Add explanatory text using manual formatting for proper bold text
        text_box = slide.shapes.add_textbox(
            Inches(text_left), 
            Inches(table_top), 
            Inches(text_width), 
            Inches(text_height)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # Add explanations for each metric with proper bold formatting
        first_metric = True
        for metric, explanation_text in metrics_text.items():
            if first_metric:
                # Use the existing first paragraph
                p = text_frame.paragraphs[0]
                first_metric = False
            else:
                # Add new paragraphs for subsequent metrics
                p = text_frame.add_paragraph()
                
            p.text = f"{metric}:"
            p.font.bold = True  # Proper bold formatting for PowerPoint
            p.font.size = Pt(10)
            p.space_after = Pt(3)
            
            p = text_frame.add_paragraph()
            
            # Parse text to handle multiple "{prefix}: {text}" patterns
            if ": " in explanation_text:
                # Split the text into lines and process each line
                lines = explanation_text.split('\n')
                first_line = True
                
                for line in lines:
                    if not first_line:
                        # Add line break for subsequent lines
                        p.add_run().text = "\n"
                    first_line = False
                    
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Check if this line has the "{prefix}: {text}" pattern
                    if ": " in line and not line.startswith(" "):  # Avoid matching indented continuations
                        prefix, rest_text = line.split(": ", 1)
                        
                        # Guardrail: only format as prefix if it's a reasonable length (≤12 characters)
                        if len(prefix) <= 12:
                            # Add prefix with bold and underline
                            run = p.add_run()
                            run.text = prefix + ": "
                            run.font.bold = True
                            run.font.underline = True
                            run.font.size = Pt(10)
                            
                            # Add rest of text with normal formatting
                            run = p.add_run()
                            run.text = rest_text
                            run.font.size = Pt(10)
                        else:
                            # Prefix too long, treat as normal text
                            run = p.add_run()
                            run.text = line
                            run.font.size = Pt(10)
                    else:
                        # No colon pattern, use normal formatting
                        run = p.add_run()
                        run.text = line
                        run.font.size = Pt(10)
            else:
                # No colon found anywhere, use normal formatting
                p.text = explanation_text
                p.font.size = Pt(10)
            
            p.space_after = Pt(10)
    
    def _calculate_content_layout(self, content_type: str, text_height: float = 0, 
                                table_width: float = 0, table_height: float = 0,
                                figure_width: float = 0, figure_height: float = 0) -> Dict[str, Dict[str, float]]:
        """Calculate optimal layout positioning for different content types."""
        # Validate layout type
        content_type = self._validate_layout_type(content_type)
        
        layouts = {}
        
        # Handle 'auto' layout by determining the best layout based on content
        if content_type == 'auto':
            # Automatically choose the best layout based on available content
            has_text = text_height > 0
            has_tables = table_width > 0 and table_height > 0
            has_figure = figure_width > 0 and figure_height > 0
            
            if has_figure and has_text and has_tables:
                content_type = 'text_tables_figure'
            elif has_figure and has_text:
                content_type = 'text_figure'
            elif has_text and has_tables:
                content_type = 'text_tables'
            elif has_figure:
                content_type = 'figure'
            else:
                content_type = 'text'  # Default to text layout
        
        if content_type == 'text':
            layouts['text'] = {
                'left': CONTENT_LEFT,
                'top': CONTENT_TOP,
                'width': CONTENT_WIDTH,
                'height': CONTENT_HEIGHT
            }
            
        elif content_type == 'text_figure':
            # Determine layout based on text size
            if text_height < CONTENT_HEIGHT * 0.25:
                # Text on top, figure below
                layouts['text'] = {
                    'left': CONTENT_LEFT,
                    'top': CONTENT_TOP,
                    'width': CONTENT_WIDTH,
                    'height': text_height
                }
                layouts['figure'] = {
                    'left': CONTENT_LEFT,
                    'top': CONTENT_TOP + text_height + ELEMENT_GAP,
                    'width': CONTENT_WIDTH,
                    'height': CONTENT_HEIGHT - text_height - ELEMENT_GAP
                }
            else:
                # Side by side layout
                text_width = min(CONTENT_WIDTH * 0.4, 6.0)
                layouts['text'] = {
                    'left': CONTENT_LEFT,
                    'top': CONTENT_TOP,
                    'width': text_width,
                    'height': CONTENT_HEIGHT
                }
                layouts['figure'] = {
                    'left': CONTENT_LEFT + text_width + ELEMENT_GAP,
                    'top': CONTENT_TOP,
                    'width': CONTENT_WIDTH - text_width - ELEMENT_GAP,
                    'height': CONTENT_HEIGHT
                }
                
        elif content_type == 'text_tables':
            layouts['text'] = {
                'left': CONTENT_LEFT,
                'top': CONTENT_TOP,
                'width': CONTENT_WIDTH,
                'height': text_height
            }
            layouts['tables'] = {
                'left': CONTENT_LEFT,
                'top': CONTENT_TOP + text_height + ELEMENT_GAP,
                'width': CONTENT_WIDTH,
                'height': CONTENT_HEIGHT - text_height - ELEMENT_GAP
            }
            
        elif content_type == 'text_tables_figure':
            layouts['text'] = {
                'left': CONTENT_LEFT,
                'top': CONTENT_TOP,
                'width': CONTENT_WIDTH,
                'height': text_height
            }
            remaining_top = CONTENT_TOP + text_height + ELEMENT_GAP
            remaining_height = CONTENT_HEIGHT - text_height - ELEMENT_GAP
            
            layouts['tables'] = {
                'left': CONTENT_LEFT,
                'top': remaining_top,
                'width': min(table_width, CONTENT_WIDTH * 0.4),
                'height': min(table_height, remaining_height)
            }
            layouts['figure'] = {
                'left': CONTENT_LEFT + table_width + ELEMENT_GAP,
                'top': remaining_top,
                'width': CONTENT_WIDTH - table_width - ELEMENT_GAP,
                'height': remaining_height
            }
            
        return layouts
    
    def _create_basic_slide_structure(self, title: str) -> object:
        """Create a basic slide with title."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_title(slide, title)
        return slide
    
    def _add_content_block(self, slide: object, content_type: str, content_data: Dict, 
                          layout_info: Dict[str, float]) -> None:
        """Add a content block (text, table, or figure) to a slide at specified position."""
        if content_type == 'text':
            text_content = content_data.get('text', '')
            if text_content:
                self._add_text(
                    slide=slide,
                    text=text_content,
                    left=layout_info['left'],
                    top=layout_info['top'],
                    width=layout_info['width'],
                    height=layout_info['height']
                )
                
        elif content_type == 'table':
            df = content_data.get('df')
            highlight_cells = content_data.get('highlight_cells', {})
            if df is not None:
                self._add_table(
                    slide=slide,
                    df=df,
                    left=layout_info['left'],
                    top=layout_info['top'],
                    width=layout_info['width'],
                    height=layout_info['height'],
                    highlight_cells=highlight_cells
                )
                
        elif content_type == 'figure':
            figure_path = content_data.get('figure_path')
            min_size = content_data.get('min_size')
            caption = content_data.get('caption')
            if figure_path:
                self._add_figure(
                    slide=slide,
                    figure_path=figure_path,
                    left=layout_info['left'],
                    top=layout_info['top'],
                    width=layout_info.get('width'),
                    height=layout_info.get('height'),
                    min_size=min_size,
                    caption=caption
                )
    
    def _combine_text_chunks(self, text_chunks: List[str]) -> str:
        """Combine text chunks with proper formatting."""
        return '\n\n'.join(chunk.strip() for chunk in text_chunks if chunk.strip())

    def _should_use_two_column_layout(self, text_chunks: List[str], table_positions: List[Tuple[str, pd.DataFrame]]) -> bool:
        """Determine if a two-column layout should be used for text-tables."""
        # Use two-column layout when we have:
        # 1. Multiple tables (2 or more)
        # 2. At least some text chunks
        # 3. Tables are reasonably sized (not too wide)
        
        if len(table_positions) < 2 or len(text_chunks) == 0:
            return False
        
        # Calculate column width
        column_width = (CONTENT_WIDTH - ELEMENT_GAP) / 2
        
        # Check if tables are reasonably sized for two-column layout
        # Use column width as max_width when testing suitability
        max_table_width = 0
        for _, df in table_positions:
            table_width, _ = self._calculate_table_dimensions(df, max_width=column_width * 0.9)
            max_table_width = max(max_table_width, table_width)
        
        # Only use two-column if tables can fit in the column
        return max_table_width <= column_width * 0.9  # Leave some margin
    
    def _create_two_column_text_tables_layout(self, slide: object, text_chunks: List[str], 
                                             table_positions: List[Tuple[str, pd.DataFrame]], 
                                             highlight_cells: Optional[Dict]):
        """Create a two-column layout with text descriptions paired with their corresponding tables."""
        # Calculate column dimensions
        column_width = (CONTENT_WIDTH - ELEMENT_GAP) / 2
        left_column_left = CONTENT_LEFT
        right_column_left = CONTENT_LEFT + column_width + ELEMENT_GAP
        
        # Pair text chunks with tables intelligently
        pairs = self._pair_text_with_tables(text_chunks, table_positions)
        
        # Position pairs in two columns
        left_column_top = CONTENT_TOP
        right_column_top = CONTENT_TOP
        
        for i, (text_chunk, table_info) in enumerate(pairs):
            # Determine which column to use (alternate)
            use_left_column = (i % 2 == 0)
            
            if use_left_column:
                column_left = left_column_left
                current_top = left_column_top
            else:
                column_left = right_column_left
                current_top = right_column_top
            
            # Add text description
            if text_chunk:
                text_height = self._estimate_text_height(text_chunk)
                self._add_text(
                    slide=slide,
                    text=text_chunk,
                    left=column_left,
                    top=current_top,
                    width=column_width,
                    height=text_height
                )
                current_top += text_height + ELEMENT_GAP * 0.5  # Small gap between text and table
            
            # Add table
            if table_info:
                table_key, df = table_info
                table_width, table_height = self._calculate_table_dimensions(df)
                
                self._add_table(
                    slide=slide,
                    df=df,
                    left=column_left,
                    top=current_top,
                    width=min(table_width, column_width),
                    height=table_height,
                    highlight_cells=highlight_cells
                )
                current_top += table_height + ELEMENT_GAP
            
            # Update column top for next iteration
            if use_left_column:
                left_column_top = current_top
            else:
                right_column_top = current_top
    
    def _pair_text_with_tables(self, text_chunks: List[str], 
                              table_positions: List[Tuple[str, pd.DataFrame]]) -> List[Tuple[Optional[str], Optional[Tuple[str, pd.DataFrame]]]]:
        """Intelligently pair text chunks with their corresponding tables."""
        pairs = []
        
        # Strategy: 
        # - If we have equal or more text chunks than tables, pair them 1:1
        # - If we have more tables than text chunks, some tables won't have descriptions
        # - If we have fewer tables than text chunks, some text will be grouped
        
        num_text = len(text_chunks)
        num_tables = len(table_positions)
        
        if num_text >= num_tables:
            # Each table gets a text description (possibly combining multiple text chunks)
            for i in range(num_tables):
                text_chunk = text_chunks[i] if i < num_text else None
                table_info = table_positions[i]
                pairs.append((text_chunk, table_info))
            
            # Add any remaining text chunks as standalone text items
            for i in range(num_tables, num_text):
                pairs.append((text_chunks[i], None))
        else:
            # More tables than text - some tables won't have descriptions
            for i in range(max(num_text, num_tables)):
                text_chunk = text_chunks[i] if i < num_text else None
                table_info = table_positions[i] if i < num_tables else None
                if text_chunk or table_info:  # Only add if we have content
                    pairs.append((text_chunk, table_info))
        
        return pairs


# The decorator 
def dual_output(console: bool = True, slide: bool = True, 
               slide_builder: Optional[SlideLayouts] = None,
               layout_type: str = 'auto',
               show_figures: bool = True):
    """
    Decorator for functions that generate both console output and slides.
    
    Args:
        console: Whether to display console output
        slide: Whether to add to slide presentation
        slide_builder: SlideLayouts instance to add slides to
        layout_type: Layout type for slides ('auto', 'text', 'text_figure', 'text_tables', 'text_tables_figure', 'text_tables_then_figure', 'figure')
        show_figures: Whether to display matplotlib figures
    """
    def decorator(func):
        def wrapper(*args, **kwargs):
            # Validate layout type early
            if slide and slide_builder:
                slide_builder._validate_layout_type(layout_type)
            
            # Execute the function to get SlideContent
            content = func(*args, **kwargs)
            
            if console:
                print(content.render_console())
            
            if slide and slide_builder:
                slide_builder.add_content_slide(content, layout_type=layout_type)
            
            return content, content.get_all_dataframes()
        return wrapper
    return decorator


# Example figure generators
def create_bar_chart(df: pd.DataFrame, **params):
    """Example figure generator - bar chart."""
    fig, ax = plt.subplots(figsize=(8, 5))
    df.plot(kind='bar', ax=ax)
    ax.set_title(params.get('chart_title', 'Data Overview'))
    ax.set_ylabel(params.get('ylabel', 'Values'))
    plt.xticks(rotation=45)
    plt.tight_layout()
    return fig


def create_scatter_plot(df: pd.DataFrame, **params):
    """Example figure generator - scatter plot."""
    fig, ax = plt.subplots(figsize=(8, 5))
    
    if len(df.columns) >= 2:
        x_col, y_col = df.columns[:2]
        ax.scatter(df[x_col], df[y_col])
        ax.set_xlabel(x_col)
        ax.set_ylabel(y_col)
        ax.set_title(params.get('chart_title', f'{y_col} vs {x_col}'))
    
    plt.tight_layout()
    return fig


if __name__ == "__main__":
    # Example usage
    data = {
        'Conversion Rate': [0.12, 0.08, 0.11, 0.13],
        'AOV': [75.0, 65.0, 80.0, 85.0]
    }
    df = pd.DataFrame(data, index=['Global', 'NA', 'EU', 'Asia'])
    
    # Create slide builder
    slides = SlideLayouts()
    
    # Example 1: Text + Table with decorator (no figures)
    @dual_output(console=True, slide=True, slide_builder=slides, layout_type='text_tables')
    def create_summary_analysis():
        return SlideContent(
            title="Summary Analysis",
            text_template="Key insight: {{ region }} shows {{ trend }} performance with {{ metric }} being {{ pct }}% {{ direction }}.\n\n{{ main_table }}\n\nThis requires immediate attention.",
            dfs={'main_table': df},
            template_params={
                'region': 'North America',
                'trend': 'concerning',
                'metric': 'conversion rate',
                'pct': 33,
                'direction': 'lower'
            }
        )
    
    # Example 2: Modern figure generation with figure_generators approach
    @dual_output(console=True, slide=True, slide_builder=slides, layout_type='text_figure')
    def create_chart_analysis():
        return SlideContent(
            title="Performance by Region",
            text_template="The chart shows {{ insight }}. {{ region }} is an outlier.",
            dfs={},  # No tables for this example
            template_params={
                'insight': 'significant regional variation',
                'region': 'North America'
            },
            figure_generators=[
                {
                    'title_suffix': 'Bar Chart',
                    'params': {
                        'df': df,
                        'chart_title': 'Regional Performance Comparison',
                        'ylabel': 'Values'
                    },
                    'function': create_bar_chart
                }
            ]
        )
    
    # Example 3: Multiple figures - each gets its own slide automatically
    @dual_output(console=True, slide=True, slide_builder=slides, layout_type='text_figure')
    def create_multi_chart_analysis():
        return SlideContent(
            title="Multi-Chart Analysis",
            text_template="Multiple perspectives on the data: {{ insight }}.",
            dfs={},
            template_params={
                'insight': 'both bar and scatter views reveal patterns'
            },
            figure_generators=[
                {
                    'title_suffix': 'Bar View',
                    'params': {
                        'df': df,
                        'chart_title': 'Regional Performance',
                        'ylabel': 'Values'
                    },
                    'function': create_bar_chart
                },
                {
                    'title_suffix': 'Scatter View',
                    'params': {
                        'df': df,
                        'chart_title': 'Regional Performance'
                    },
                    'function': create_scatter_plot
                }
            ]
        )
    
    # Run examples
    print("MODERN DECORATOR APPROACH WITH FIGURE GENERATION")
    print("=" * 60)
    
    print("\n📊 SUMMARY ANALYSIS (No Figures):")
    content1, results1 = create_summary_analysis()
    
    print("\n📈 SINGLE CHART ANALYSIS (New figure_generators approach):")
    content2, results2 = create_chart_analysis()
    
    print("\n🎯 MULTI-CHART ANALYSIS (Multiple slides auto-generated):")
    content3, results3 = create_multi_chart_analysis()
    
    # Save slides
    filepath = slides.save("modern_decorator_demo", "../output")
    print(f"\n✅ Slides saved: {filepath}")
    print("✅ Figure files automatically cleaned up!")
    print("\nKEY POINTS:")
    print("• Use 'figure_generators' parameter (not in template_params)")
    print("• Each figure_generator creates one slide automatically")
    print("• Store functions directly with 'function' key in figure_generators")
    print("• No redundant function storage needed!") 