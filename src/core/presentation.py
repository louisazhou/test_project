import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from typing import Optional

# Setup logging
logger = logging.getLogger(__name__)
# Configure root logger if necessary, or rely on calling script's config
# logging.basicConfig(level=logging.INFO)

# Define scopes for Google APIs
SCOPES = [
    'https://www.googleapis.com/auth/presentations',
    'https://www.googleapis.com/auth/drive.file'
]

def get_credentials():
    """Gets user credentials for Google APIs, handling token refresh/creation."""
    creds = None
    # Assume token.json and credentials.json are in the ROOT directory of the project
    # Adjust path relative to this utils.py file (src/core/utils.py)
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    token_path = os.path.join(project_root, 'token.json')
    credentials_path = os.path.join(project_root, 'credentials.json')

    if not os.path.exists(credentials_path):
         print(f"ERROR: credentials.json not found at {credentials_path}")
         print("Please download credentials from Google Cloud Console and place it in the project root.")
         return None

    if os.path.exists(token_path):
        try:
            creds = Credentials.from_authorized_user_file(token_path, SCOPES)
        except Exception as e:
            print(f"Error loading token.json: {e}. Will attempt re-authentication.")
            creds = None # Force re-auth

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(Request())
            except Exception as e:
                print(f"Error refreshing token: {e}. Please re-authenticate.")
                # Optionally delete token.json to force full re-auth
                # if os.path.exists(token_path): os.remove(token_path)
                flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                creds = flow.run_local_server(port=0)
        else:
            try:
                flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                creds = flow.run_local_server(port=0)
            except FileNotFoundError:
                 print(f"ERROR: credentials.json not found at {credentials_path}")
                 return None
            except Exception as e:
                 print(f"Error during authentication flow: {e}")
                 return None
        try:
            with open(token_path, 'w') as token:
                token.write(creds.to_json())
        except Exception as e:
            print(f"Error saving token to {token_path}: {e}")

    return creds

def generate_ppt(analysis_results: dict, output_dir: str, ppt_filename: str = "RCA_Summary.pptx", report_format: str = "detailed"):
    """Generates a PowerPoint presentation from RCA results and visualizations.

    Args:
        analysis_results: Dictionary containing processed analysis results per metric,
                          structured like {metric_name: {'primary_region': ..., ...}}.
        output_dir: Directory where visualization PNG files are saved and
                    where the PowerPoint file will be saved.
        ppt_filename: The name for the output PowerPoint file.
        report_format: Type of report/visualization expected ("detailed", "succinct", "individual").
    """
    prs = Presentation()
    # Use a blank slide layout (ensure layout index 6 is indeed blank in default template)
    try:
        blank_slide_layout = prs.slide_layouts[6]
    except IndexError:
        logger.warning("Blank slide layout (index 6) not found. Using first layout.")
        blank_slide_layout = prs.slide_layouts[0]

    logger.info(f"Generating PowerPoint presentation: {ppt_filename} using {report_format} format preference")

    metrics_processed = 0
    for metric, results in analysis_results.items():
        primary_region = results.get('primary_region', 'Unknown')

        # Determine the expected summary image filename based on report_format
        img_filename = None
        if report_format in ["detailed", "succinct"] and primary_region not in ["NoAnomaly", "NoData", None]:
             # TODO: Need the actual summary plot generation logic first
             # This filename assumes a summary plot is generated with this convention
             img_filename = f"{metric}_{primary_region}_{report_format}_summary.png"
        elif report_format == "individual" or primary_region in ["NoAnomaly", "NoData", None]:
             # Fallback to the individual metric plot
             img_filename = f"metric_{metric}.png"
        # Add more specific fallbacks if needed (e.g., anomaly_only plot)

        img_path = None
        if img_filename:
            candidate_path = os.path.join(output_dir, img_filename)
            if os.path.exists(candidate_path):
                img_path = candidate_path
                logger.info(f"Found image file for metric '{metric}': {img_filename}")
            else:
                 logger.warning(f"Expected image file not found for metric '{metric}': {candidate_path}")
                 # Try falling back to basic metric plot if summary wasn't found
                 fallback_filename = f"metric_{metric}.png"
                 fallback_path = os.path.join(output_dir, fallback_filename)
                 if os.path.exists(fallback_path):
                     img_path = fallback_path
                     logger.info(f"Using fallback image file: {fallback_filename}")
                 else:
                      logger.warning(f"Fallback image file also not found: {fallback_path}")
        else:
            logger.warning(f"Could not determine image filename for metric '{metric}' with format '{report_format}'.")


        if img_path:
            logger.info(f"Adding slide for metric: {metric}")
            slide = prs.slides.add_slide(blank_slide_layout)

            # Adjust layout based on image type. Standard slide is 10x7.5 inches
            # These dimensions assume landscape summary plots; adjust if needed
            img_width_inches = 9.8
            left = Inches(0.1)
            top = Inches(0.1)

            try:
                slide.shapes.add_picture(img_path, left, top, width=Inches(img_width_inches))
                metrics_processed += 1
            except Exception as e:
                logger.error(f"Error adding image {img_path} to slide: {e}")
        # else: (logging already handled above)

    # Save the presentation
    if metrics_processed > 0:
        ppt_path = os.path.join(output_dir, ppt_filename)
        try:
            prs.save(ppt_path)
            logger.info(f"PowerPoint presentation saved successfully to {ppt_path}")
            return ppt_path
        except Exception as e:
            logger.error(f"Error saving PowerPoint presentation to {ppt_path}: {e}")
            return None # Return None if save fails
    else:
        logger.warning("No metric images were found or added to the PowerPoint presentation.")
        return None # Return None if no slides added

def upload_to_drive(file_path: str, folder_id: Optional[str] = None, mime_type: str = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'):
    """Upload a file (e.g., the generated PPTX) to Google Drive.

    Args:
        file_path: Path to the file to upload (e.g., the PPTX path).
        folder_id: Optional Google Drive folder ID to upload into.
        mime_type: MIME type of the file being uploaded.

    Returns:
        File ID of the uploaded file if successful, None otherwise.
    """
    if not file_path or not os.path.exists(file_path):
        logger.error(f"File not found or path is invalid: {file_path}")
        return None

    try:
        creds = get_credentials()
        if not creds:
            logger.error("Failed to get Google credentials. Cannot upload.")
            return None

        service = build('drive', 'v3', credentials=creds)
        file_name = os.path.basename(file_path)

        # Metadata for the file on Google Drive
        file_metadata = {
            'name': file_name,
            # Optionally convert to Google Slides format on upload
            'mimeType': 'application/vnd.google-apps.presentation'
        }
        if folder_id:
            file_metadata['parents'] = [folder_id]

        # Prepare the media for upload
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)

        # Perform the upload
        file = service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, webViewLink' # Request fields needed
        ).execute()

        file_id = file.get('id')
        web_link = file.get('webViewLink')

        logger.info(f"Successfully uploaded '{file_name}' to Google Drive.")
        logger.info(f"File ID: {file_id}")
        logger.info(f"Web Link: {web_link}")

        return file_id

    except Exception as e:
        logger.error(f"Error uploading file '{file_path}' to Google Drive: {e}", exc_info=True)
        return None 