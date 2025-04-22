import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from utils import get_credentials

# Setup logging
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

def generate_ppt(analysis_results: dict, output_dir: str, ppt_filename: str = "RCA_Summary.pptx", visualization_type: str = "detailed"):
    """Generates a PowerPoint presentation from RCA visualization images.

    Args:
        analysis_results: Dictionary containing processed analysis results per metric,
                          structured like {metric_name: {'primary_region': ..., ...}}.
                          Used to find the correct image files.
        output_dir: Directory where the visualization PNG files are saved and
                    where the PowerPoint file will be saved.
        ppt_filename: The name for the output PowerPoint file.
        visualization_type: Either "detailed" or "succinct" - determines which image files to include.
    """
    prs = Presentation()
    # Use a blank slide layout (usually layout index 6, but can vary)
    blank_slide_layout = prs.slide_layouts[6] 
    
    logger.info(f"Generating PowerPoint presentation: {ppt_filename} using {visualization_type} visualizations")
    
    metrics_processed = 0
    for metric, results in analysis_results.items():
        primary_region = results.get('primary_region', 'Unknown')
        has_root_cause = results.get('best_hypothesis_name') is not None
        
        # List of possible image filenames to try, in order of preference
        img_filename_candidates = []
        
        # 1. First check for the requested visualization type
        if has_root_cause and primary_region not in ["NoAnomaly", "NoData", None]:
            # Look for the appropriately named image based on visualization type
            img_filename_candidates.append(f"{metric}_{primary_region}_{visualization_type}_summary.png")
        
        # 2. Then check for the alternative visualization type if first one doesn't exist
        alt_viz_type = "detailed" if visualization_type == "succinct" else "succinct"
        if has_root_cause and primary_region not in ["NoAnomaly", "NoData", None]:
            img_filename_candidates.append(f"{metric}_{primary_region}_{alt_viz_type}_summary.png")
            
        # 3. For backward compatibility, check for old naming convention
        img_filename_candidates.append(f"{metric}_{primary_region}_RCA_Summary.png")
            
        # 4. Check for anomaly-only visualizations 
        if primary_region not in ["NoAnomaly", "NoData", None]:
            img_filename_candidates.append(f"{metric}_{primary_region}_anomaly_only.png")
            img_filename_candidates.append(f"{metric}_AnomalyOnly_{visualization_type}_summary.png")
        
        # 5. Finally check for performance-only fallback
        img_filename_candidates.append(f"{metric}_performance.png")
        
        # Try to find one of the candidate images
        img_path = None
        for candidate in img_filename_candidates:
            candidate_path = os.path.join(output_dir, candidate)
            if os.path.exists(candidate_path):
                img_path = candidate_path
                logger.info(f"Found image file for metric '{metric}': {candidate}")
                break
        
        if img_path:
            logger.info(f"Adding slide for metric: {metric}")
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Adjust layout based on image type. Standard slide is 10x7.5 inches
            img_width_inches = 9.8 # Default for full-page visualizations
            left = Inches(0.1)
            top = Inches(0.1)
            
            # Adjust size for performance-only visualizations which might be smaller
            if "performance" in img_path:
                img_width_inches = 7.0
                left = Inches(1.5)
                top = Inches(1.0)
            
            try:
                pic = slide.shapes.add_picture(img_path, left, top, width=Inches(img_width_inches))
                metrics_processed += 1
            except Exception as e:
                logger.error(f"Error adding image {img_path} to slide: {e}")
        else:
            logger.warning(f"No suitable image file found for metric '{metric}'. Tried: {', '.join(img_filename_candidates)}")

    if metrics_processed > 0:
        ppt_path = os.path.join(output_dir, ppt_filename)
        try:
            prs.save(ppt_path)
            logger.info(f"PowerPoint presentation saved successfully to {ppt_path}")
            return ppt_path
        except Exception as e:
            logger.error(f"Error saving PowerPoint presentation to {ppt_path}: {e}")
    else:
        logger.warning("No metric images were found to add to the PowerPoint presentation.")
        
    return ppt_path if metrics_processed > 0 else None

def upload_to_drive(file_path: str, folder_id: str = None, mime_type: str = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'):
    """Upload a file to Google Drive.
    
    Args:
        file_path: Path to the file to upload
        folder_id: Optional Google Drive folder ID to upload to (uses Drive root if not specified)
        mime_type: MIME type of the file (default is for PPTX files)
        
    Returns:
        File ID of the uploaded file if successful, None otherwise
    """
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        return None
        
    try:
        # Get credentials and build service
        creds = get_credentials()
        service = build('drive', 'v3', credentials=creds)
        
        # File metadata
        file_name = os.path.basename(file_path)
        file_metadata = {
            'name': file_name
        }
        
        # Add to specific folder if provided
        if folder_id:
            file_metadata['parents'] = [folder_id]
            
        # Upload file
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        file = service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id,webViewLink'
        ).execute()
        
        file_id = file.get('id')
        web_link = file.get('webViewLink')
        
        logger.info(f"Successfully uploaded {file_name} to Google Drive")
        logger.info(f"File ID: {file_id}")
        logger.info(f"Web link: {web_link}")
        
        return file_id
        
    except Exception as e:
        logger.error(f"Error uploading to Google Drive: {e}")
        return None 